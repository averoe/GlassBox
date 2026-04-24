"""Multimodal utilities for handling text, images, and PDFs.

Real implementations for text chunking, PDF extraction, and
OCR-based image text extraction.
"""

from __future__ import annotations

from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, field

from glassbox_rag.utils.logging import get_logger

logger = get_logger(__name__)


@dataclass
class MultimodalContent:
    """Represents multimodal content (text, image, PDF, pptx)."""

    content_type: str  # "text", "image", "pdf", "pptx"
    content: Any  # bytes or string
    encoding: str = "utf-8"
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "content_type": self.content_type,
            "encoding": self.encoding,
            "metadata": self.metadata,
        }


class BaseMultimodalProcessor(ABC):
    """Base class for multimodal processors."""

    @abstractmethod
    async def process(self, content: MultimodalContent) -> list[str]:
        """
        Process multimodal content into text chunks.

        Args:
            content: Multimodal content to process.

        Returns:
            List of text strings extracted from content.

        Raises:
            ValueError: If content type doesn't match processor.
            ProcessingError: If extraction fails.
        """


class ProcessingError(Exception):
    """Raised when content processing fails."""


class TextProcessor(BaseMultimodalProcessor):
    """Processor for plain text content."""

    async def process(self, content: MultimodalContent) -> list[str]:
        if content.content_type != "text":
            raise ValueError(f"Expected text content, got {content.content_type}")

        text = content.content
        if isinstance(text, bytes):
            text = text.decode(content.encoding)

        if not text or not text.strip():
            return []

        return [text]


class ImageProcessor(BaseMultimodalProcessor):
    """Processor for image content using OCR."""

    async def process(self, content: MultimodalContent) -> list[str]:
        if content.content_type != "image":
            raise ValueError(f"Expected image content, got {content.content_type}")

        try:
            from PIL import Image
            import io

            image_data = content.content
            if isinstance(image_data, str):
                # File path
                image = Image.open(image_data)
            else:
                image = Image.open(io.BytesIO(image_data))

            # Try OCR with pytesseract
            try:
                import asyncio
                import pytesseract
                # pytesseract is a blocking C extension — run in executor
                loop = asyncio.get_running_loop()
                text = await loop.run_in_executor(
                    None, pytesseract.image_to_string, image
                )
                if text and text.strip():
                    return [text.strip()]
            except ImportError:
                logger.warning(
                    "pytesseract not installed, cannot extract text from images. "
                    "Run: pip install glassbox-rag[multimodal]"
                )

            # Fallback: return image metadata as context
            width, height = image.size
            mode = image.mode
            return [
                f"[Image: {width}x{height}, mode={mode}, "
                f"format={image.format or 'unknown'}]"
            ]

        except ImportError:
            raise ProcessingError(
                "Pillow not installed. Run: pip install glassbox-rag[multimodal]"
            )
        except Exception as e:
            raise ProcessingError(f"Image processing failed: {e}") from e


class PDFProcessor(BaseMultimodalProcessor):
    """Processor for PDF documents using pypdf."""

    async def process(self, content: MultimodalContent) -> list[str]:
        if content.content_type != "pdf":
            raise ValueError(f"Expected PDF content, got {content.content_type}")

        try:
            from pypdf import PdfReader
            import io

            pdf_data = content.content
            if isinstance(pdf_data, str):
                reader = PdfReader(pdf_data)
            else:
                reader = PdfReader(io.BytesIO(pdf_data))

            max_pages = content.metadata.get("max_pages", 100)
            extract_images = content.metadata.get("extract_images", False)

            texts: list[str] = []
            for i, page in enumerate(reader.pages[:max_pages]):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    texts.append(page_text.strip())

                # Image extraction from PDF pages
                if extract_images:
                    for image_obj in page.images:
                        try:
                            image_processor = ImageProcessor()
                            image_content = MultimodalContent(
                                content_type="image",
                                content=image_obj.data,
                            )
                            image_texts = await image_processor.process(image_content)
                            texts.extend(image_texts)
                        except Exception as e:
                            logger.debug("Failed to extract image from page %d: %s", i, e)

            if not texts:
                logger.warning("No text extracted from PDF (%d pages)", len(reader.pages))

            return texts

        except ImportError:
            raise ProcessingError(
                "pypdf not installed. Run: pip install glassbox-rag[multimodal]"
            )
        except Exception as e:
            raise ProcessingError(f"PDF processing failed: {e}") from e


class PPTXProcessor(BaseMultimodalProcessor):
    """Processor for PowerPoint (.pptx) files."""

    async def process(self, content: MultimodalContent) -> list[str]:
        if content.content_type != "pptx":
            raise ValueError(f"Expected pptx content, got {content.content_type}")

        try:
            from pptx import Presentation
            import io

            pptx_data = content.content
            if isinstance(pptx_data, str):
                prs = Presentation(pptx_data)
            else:
                prs = Presentation(io.BytesIO(pptx_data))

            texts: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_texts.append(para_text)

                if slide_texts:
                    texts.append(
                        f"[Slide {slide_num}]\n" + "\n".join(slide_texts)
                    )

            return texts

        except ImportError:
            raise ProcessingError(
                "python-pptx not installed. Run: pip install glassbox-rag[multimodal]"
            )
        except Exception as e:
            raise ProcessingError(f"PPTX processing failed: {e}") from e


class MultimodalHandler:
    """
    Unified handler for multimodal content.

    Provides an interface for processing text, images, PDFs,
    and PPTX in a consistent manner. Supports registering
    custom processors.
    """

    def __init__(self) -> None:
        self.processors: dict[str, BaseMultimodalProcessor] = {
            "text": TextProcessor(),
            "image": ImageProcessor(),
            "pdf": PDFProcessor(),
            "pptx": PPTXProcessor(),
        }

    def register_processor(
        self,
        content_type: str,
        processor: BaseMultimodalProcessor,
    ) -> None:
        """Register a custom processor for a content type."""
        self.processors[content_type] = processor
        logger.info("Registered multimodal processor: %s", content_type)

    async def process(self, content: MultimodalContent) -> list[str]:
        """
        Process multimodal content.

        Returns:
            List of text chunks extracted from content.

        Raises:
            ValueError: If content type is not supported.
            ProcessingError: If extraction fails.
        """
        if content.content_type not in self.processors:
            supported = list(self.processors.keys())
            raise ValueError(
                f"Unsupported content type: '{content.content_type}'. "
                f"Supported types: {supported}"
            )

        processor = self.processors[content.content_type]
        results = await processor.process(content)

        logger.debug(
            "Processed %s content: %d text chunks extracted",
            content.content_type,
            len(results),
        )
        return results

    @property
    def supported_types(self) -> list[str]:
        """List supported content types."""
        return list(self.processors.keys())
